import os
import re
import json
import tkinter as tk
from tkinter import messagebox
from typing import Any, Dict, List, Tuple

import requests
from dotenv import load_dotenv

# Optional Tk globals used by reframe_action() when this module is wired into a UI.
root = None
user_text = None
file_content = None
output_text = None

# ----------------------------
# Load Environment Variables
# ----------------------------
load_dotenv()

DATABRICKS_URL = os.getenv("DATABRICKS_URL")
DATABRICKS_API_KEY = os.getenv("DATABRICKS_API")

if not DATABRICKS_URL or not DATABRICKS_API_KEY:
    print("Warning: Databricks URL or API key is missing.")
else:
    print("Databricks environment variables loaded successfully.")

# ----------------------------
# Constants
# ----------------------------
# Applied Materials part number formats:
#   0495-12345
#   ESW0495-12345 / ESW0020-12345
PN_FRAGMENT = r"(?:ESW)?\d{4,5}-\d{5}"
PART_NUMBER_PATTERN = re.compile(rf"\b{PN_FRAGMENT}\b", re.IGNORECASE)

TOP_LEVEL_SECTION_NAMES = [
    "Title",
    "Problem Statement",
    "Solution Statement",
    "Proposed Solution",
]

PROBLEM_SUBSECTION_NAMES = [
    "Reference Change Summary",
    "Affected/Impacted Part Numbers",
    "Affected Part Numbers",
    "Issue Part Number(s)",
    "Issue Part Numbers",
    "Problem Description",
    "Impact Details",
]

SOLUTION_SUBSECTION_NAMES = [
    "Proposed Change",
    "Solution Description",
    "Benefits of the Proposed Solution",
]

FIRST_PERSON_TERMS = r"\b(I|we|our|my|us)\b"
URGENCY_TERMS = [r"\burgent\b", r"\burgently\b", r"\bimmediately\b", r"\basap\b"]
BLAME_TERMS = [r"\bblame\b", r"\bfault\b", r"\bnegligence\b", r"\bmistake by\b"]
DISALLOWED_PROBLEM_TERMS = [
    r"\bpropos(?:e|ed|al)\b",
    r"\breplace(?:d|ment)?\b",
    r"\bsolution\b",
    r"\bapproval\b",
    r"\brecommend(?:ed|ation)?\b",
]

ENGINEERING_SHORTCUTS: List[Tuple[str, str]] = [
    ("Engineering Change Order / Notice", "ECO"),
    ("Failure Mode & Effects Analysis", "FMEA"),
    ("Material Requirement Planning", "MRP"),
    ("Standard Operating Procedure", "SOP"),
    ("Engineering Change Request", "ECR"),
    ("Engineering Change Notice", "ECN"),
    ("Minimum Order Quantity", "MOQ"),
    ("Approved Vendor List", "AVL"),
    ("First Article Inspection", "FAI"),
    ("Manufacturing Drawing", "Mfg. DWG"),
    ("Root Cause Analysis", "RCA"),
    ("Purchase Requisition", "PR"),
    ("Purchase Order", "PO"),
    ("Bill of Material", "BOM"),
    ("Lead Time", "LT"),
    ("Process FMEA", "PFMEA"),
    ("Design FMEA", "DFMEA"),
    ("First In First Out", "FIFO"),
    ("Fit, Form & Function", "FFF"),
    ("Goods Receipt", "GR"),
    ("Drawing", "DWG"),
    ("Revision", "REV"),
    ("Assembly", "Assy."),
]

# ----------------------------
# Core Helpers
# ----------------------------
def extract_part_numbers(text: str) -> List[str]:
    if not text:
        return []
    seen = set()
    out: List[str] = []
    for p in PART_NUMBER_PATTERN.findall(text):
        p = p.upper()
        if p not in seen:
            seen.add(p)
            out.append(p)
    return out


def _dedupe_lines(lines: List[str]) -> List[str]:
    seen = set()
    out = []
    for ln in lines:
        key = re.sub(r"\s+", " ", ln.strip().lower())
        if key and key not in seen:
            seen.add(key)
            out.append(ln.strip())
    return out


def extract_bom_relationships(text: str) -> List[Tuple[str, str]]:
    """
    Returns (parent, child). Rule: Used in Part Number = Parent.
    """
    if not text:
        return []

    rels: List[Tuple[str, str]] = []
    seen: set[Tuple[str, str]] = set()

    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue

        parent = None
        child = None

        p_match = re.search(
            rf"parent(?:\s*part(?:\s*number)?)?\s*[:=\-]?\s*({PN_FRAGMENT})",
            line, flags=re.IGNORECASE
        )
        c_match = re.search(
            rf"child(?:\s*part(?:\s*number)?)?\s*[:=\-]?\s*({PN_FRAGMENT})",
            line, flags=re.IGNORECASE
        )
        if p_match and c_match:
            parent = p_match.group(1).upper()
            child = c_match.group(1).upper()

        if not (parent and child):
            m = re.search(
                rf"({PN_FRAGMENT}).*?(?:used\s+in|where\s+used).*?({PN_FRAGMENT})",
                line, flags=re.IGNORECASE
            )
            if m:
                child = m.group(1).upper()
                parent = m.group(2).upper()

        if not (parent and child):
            m = re.search(
                rf"used\s*in\s*part\s*number\s*[:=\-]?\s*({PN_FRAGMENT})",
                line, flags=re.IGNORECASE
            )
            if m:
                parent = m.group(1).upper()
                pns = extract_part_numbers(line)
                for pn in pns:
                    if pn != parent:
                        child = pn
                        break

        if parent and child:
            pair = (parent, child)
            if pair not in seen:
                seen.add(pair)
                rels.append(pair)

    return rels


def build_bom_context(text: str) -> str:
    rels = extract_bom_relationships(text)
    if not rels:
        return ""
    lines = ["Detected BOM context (Used in Part Number = Parent):"]
    for i, (parent, child) in enumerate(rels, start=1):
        lines.append(f"{i}. Parent Part Number: {parent}; Child Part Number: {child}; Used in Part Number: {parent}")
    return "\n".join(lines)


def extract_replacement_pairs(text: str) -> List[Tuple[str, str]]:
    """
    Returns (obsolete_part, replacement_part) pairs inferred from text.
    """
    if not text:
        return []

    pairs: List[Tuple[str, str]] = []
    seen: set[Tuple[str, str]] = set()

    patterns = [
        # "<old> replaced by <new>"
        rf"({PN_FRAGMENT})\s*(?:is\s*)?(?:to\s*be\s*)?replaced\s*by\s*({PN_FRAGMENT})",
        # "replace <old> with <new>"
        rf"replace\s*({PN_FRAGMENT})\s*with\s*({PN_FRAGMENT})",
        # "obsolete <old> ... replacement <new>"
        rf"obsolete(?:d)?\s*[:=\-]?\s*({PN_FRAGMENT}).*?replacement\s*[:=\-]?\s*({PN_FRAGMENT})",
        # "from <old> to <new>"
        rf"from\s*[:=\-]?\s*({PN_FRAGMENT})\s*to\s*[:=\-]?\s*({PN_FRAGMENT})",
    ]

    for raw_line in str(text).splitlines():
        line = raw_line.strip()
        if not line:
            continue
        for pat in patterns:
            m = re.search(pat, line, flags=re.IGNORECASE)
            if not m:
                continue
            old_pn = m.group(1).upper()
            new_pn = m.group(2).upper()
            if old_pn != new_pn:
                pair = (old_pn, new_pn)
                if pair not in seen:
                    seen.add(pair)
                    pairs.append(pair)

    return pairs


def build_replacement_context(text: str) -> str:
    pairs = extract_replacement_pairs(text)
    if not pairs:
        return ""
    lines = ["Detected replacement context (Obsolete -> Replacement):"]
    for i, (old_pn, new_pn) in enumerate(pairs, start=1):
        lines.append(f"{i}. Obsolete Part Number: {old_pn}; Replacement Part Number: {new_pn}")
    return "\n".join(lines)


def clean_ai_output(text: str) -> str:
    if not text:
        return ""
    # Remove markdown code fences that may wrap the entire output
    text = re.sub(r"^```[a-zA-Z]*\n?", "", text, flags=re.MULTILINE)
    text = re.sub(r"^```\s*$", "", text, flags=re.MULTILINE)
    text = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", text)
    text = re.sub(r"<img[^>]*>", "", text, flags=re.IGNORECASE)
    # Strip markdown heading markers (##, ###, etc.) from section headers
    text = re.sub(r"^(\s*)#{1,6}\s+", r"\1", text, flags=re.MULTILINE)
    # Strip bold/italic markers around section header names
    # e.g. **Title:** -> Title:  or **Title**:  -> Title:
    text = re.sub(r"\*{1,3}([^*\n]+?)\*{1,3}", r"\1", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def remove_unwanted_sections(text: str) -> str:
    blocked = ["Part Numbers Identified", "BOM Relationships Identified", "Quality Gate Warnings"]
    for header in blocked:
        pattern = rf"\n*{re.escape(header)}\s*:\s*(.*?)(?=\n[A-Za-z][A-Za-z ()/\-]*:\s*|\Z)"
        text = re.sub(pattern, "", text, flags=re.IGNORECASE | re.DOTALL)
    return text.strip()


def _extract_section(text: str, section_name: str) -> str:
    others = [h for h in TOP_LEVEL_SECTION_NAMES if h.lower() != section_name.lower()]
    boundary = "|".join(re.escape(h) for h in others)
    pattern = rf"^\s*{re.escape(section_name)}\s*:\s*(.*?)(?=^\s*(?:{boundary})\s*:|\Z)"
    m = re.search(pattern, text, flags=re.IGNORECASE | re.MULTILINE | re.DOTALL)
    return m.group(1).strip() if m else ""


def _extract_subsection(section_text: str, subsection_name: str, valid_names: List[str]) -> str:
    others = [n for n in valid_names if n.lower() != subsection_name.lower()]
    boundary = "|".join(re.escape(n) for n in others)
    pattern = rf"^\s*{re.escape(subsection_name)}\s*:\s*(.*?)(?=^\s*(?:{boundary})\s*:|\Z)"
    m = re.search(pattern, section_text, flags=re.IGNORECASE | re.MULTILINE | re.DOTALL)
    return m.group(1).strip() if m else ""


def _extract_title(text: str) -> str:
    m = re.search(r"^\s*Title\s*:\s*(.+)$", text, flags=re.IGNORECASE | re.MULTILINE)
    return m.group(1).strip() if m else ""


def _has_bullet_lines(text: str) -> bool:
    return bool(re.search(r"^\s*[-*•]\s+\S+", text or "", flags=re.MULTILINE))


def _compress_repeated_part_number_labels(text: str) -> str:
    """
    Converts:
      Part Number: 0150-43377 and Part Number: 0150-43378
    into:
      Part Number: 0150-43377, 0150-43378
    """
    if not text:
        return text

    pattern = rf"Part\s*Number\s*:\s*({PN_FRAGMENT})\s*(?:,|and|/)\s*Part\s*Number\s*:\s*({PN_FRAGMENT})"
    prev = None
    cur = text
    while prev != cur:
        prev = cur
        cur = re.sub(pattern, r"Part Number: \1, \2", cur, flags=re.IGNORECASE)
    return cur


def _short_crisp_bullets(text: str, max_bullets: int = 3, max_words: int = 12) -> str:
    if not text:
        return ""
    bullets: List[str] = []
    for raw in text.splitlines():
        s = raw.strip()
        if not s:
            continue
        s = re.sub(r"^[-*•]\s*", "", s).strip()
        words = s.split()
        if len(words) > max_words:
            s = " ".join(words[:max_words]).rstrip(",.;:") + "."
        bullets.append(f"- {s}")
        if len(bullets) >= max_bullets:
            break
    return "\n".join(_dedupe_lines(bullets)).strip()


def _normalize_part_scope_lines(text: str) -> str:
    """
    For Affected Part Numbers / Issue Part Number(s):
    - no repeated 'Part Number:' labels
    - no descriptions, only part numbers
    - if multiple PNs in one line, keep as comma-separated PN list only
    """
    if not text:
        return ""

    out: List[str] = []
    for raw in text.splitlines():
        line = re.sub(r"^[-*•]\s*", "", raw.strip())
        if not line:
            continue

        pns = extract_part_numbers(line)
        if not pns:
            # Skip non-AM identifiers (e.g., SPS/PCR IDs or supplier part numbers).
            continue

        uniq = []
        seen = set()
        for pn in pns:
            if pn not in seen:
                seen.add(pn)
                uniq.append(pn)

        out.append(", ".join(uniq))

    return "\n".join(_dedupe_lines(out)).strip()


def _normalize_multiline_list(text: str) -> str:
    if not text:
        return ""
    if extract_part_numbers(text):
        return _normalize_part_scope_lines(text)

    lines = []
    for raw in text.splitlines():
        ln = re.sub(r"^[-*•]\s*", "", raw.strip())
        if ln:
            lines.append(ln)
    return "\n".join(_dedupe_lines(lines)).strip()


def _linebreak_after_fullstop(text: str) -> str:
    """Put each sentence on a new line using full stop as separator."""
    if not text:
        return ""

    normalized = re.sub(r"\s+", " ", str(text).strip())
    if not normalized:
        return ""

    parts = [p.strip() for p in re.split(r"(?<=\.)\s+", normalized) if p.strip()]
    return "\n".join(parts).strip()


def _extract_revision_number(text: str) -> int | None:
    """Extract revision number from text like Rev 01 / REV.02 / Rev-3."""
    if not text:
        return None
    m = re.search(r"\brev\.?\s*[-:]?\s*0*(\d+)\b", str(text), flags=re.IGNORECASE)
    if not m:
        return None
    try:
        return int(m.group(1))
    except ValueError:
        return None


def _ensure_from_to_revision_order(from_val: str, to_val: str) -> tuple[str, str]:
    """Ensure revision transitions are old->new (e.g., Rev 01 -> Rev 02)."""
    from_rev = _extract_revision_number(from_val)
    to_rev = _extract_revision_number(to_val)

    # Only auto-swap when both sides clearly contain revision numbers.
    if from_rev is not None and to_rev is not None and from_rev > to_rev:
        return to_val, from_val
    return from_val, to_val


def _extract_revision_sequence_from_sps(payload: Dict[str, Any]) -> tuple[int, int] | None:
    """Extract older and newer revision numbers from SPS records to establish From/To order."""
    sps_records = (payload or {}).get("sps_records") or []
    if not sps_records:
        return None
    
    revisions = []
    for rec in sps_records:
        problem = str(rec.get("problem") or "")
        solution = str(rec.get("solution") or "")
        combined = f"{problem} {solution}"
        
        for match in re.finditer(r"\brev\.?\s*[-:]?\s*0*(\d+)", combined, flags=re.IGNORECASE):
            try:
                rev_num = int(match.group(1))
                if rev_num not in revisions:
                    revisions.append(rev_num)
            except ValueError:
                continue
    
    if len(revisions) >= 2:
        revisions.sort()
        return (revisions[0], revisions[-1])
    return None


def _apply_source_based_from_to_order(text: str, payload: Dict[str, Any]) -> str:
    """Apply From/To ordering based on SPS revision sequence if available."""
    if not text or not payload:
        return text
    
    rev_range = _extract_revision_sequence_from_sps(payload)
    if not rev_range:
        return text
    
    old_rev, new_rev = rev_range
    
    # Find all From: ... To: ... patterns and check if they need swapping
    def check_and_swap(m: re.Match) -> str:
        indent = m.group(1)
        from_val = m.group(2).strip()
        to_val = m.group(3).strip()
        
        from_rev = _extract_revision_number(from_val)
        to_rev = _extract_revision_number(to_val)
        
        # If we detected old_rev and new_rev from source, enforce that order
        if from_rev == new_rev and to_rev == old_rev:
            from_val, to_val = to_val, from_val
        
        return f"{indent}From: {from_val}\n{indent}To   : {to_val}"
    
    text = re.sub(
        r'^(\s*)From\s*:\s*(.+?)\s*To\s*:\s*(.+?)\s*$',
        check_and_swap,
        text,
        flags=re.IGNORECASE | re.MULTILINE
    )
    return text


def _format_from_to_changes(text: str) -> str:
    """
    Normalizes From/To change descriptions to a clean two-line format:
        From: <previous value>
        To   : <new value>

    Handles:
      - "From: X, To: Y" or "From: X; To: Y" on one line  → split to two lines
      - "From: X\nTo: Y" (no spacing)                      → align "To   :"
    """
    if not text:
        return text

    result_lines = []
    for line in text.splitlines():
        # Case 1: "From: X, To: Y" on one line — split them
        m1 = re.match(
            r'^(\s*)From\s*:\s*(.+?)\s*[,;]\s*To\s*:\s*(.+?)\s*$',
            line, flags=re.IGNORECASE
        )
        if m1:
            indent = m1.group(1)
            from_val, to_val = _ensure_from_to_revision_order(m1.group(2).strip(), m1.group(3).strip())
            result_lines.append(f"{indent}From: {from_val}")
            result_lines.append(f"{indent}To   : {to_val}")
            continue

        # Case 2: standalone "To: X" immediately after a "From:" line — normalise alignment
        m2 = re.match(r'^(\s*)To\s*:\s*(.+?)\s*$', line, flags=re.IGNORECASE)
        if m2 and result_lines and re.match(r'^\s*From\s*:', result_lines[-1], flags=re.IGNORECASE):
            prev_from = re.match(r'^(\s*)From\s*:\s*(.+?)\s*$', result_lines[-1], flags=re.IGNORECASE)
            if prev_from:
                indent = prev_from.group(1)
                from_val, to_val = _ensure_from_to_revision_order(prev_from.group(2).strip(), m2.group(2).strip())
                result_lines[-1] = f"{indent}From: {from_val}"
                result_lines.append(f"{indent}To   : {to_val}")
            else:
                indent = m2.group(1)
                result_lines.append(f"{indent}To   : {m2.group(2).strip()}")
            continue

        result_lines.append(line)

    return "\n".join(result_lines)


def _inline_problem_description_from_to(text: str) -> str:
    """
    Ensures From/To change text is part of the Problem Description sentence,
    not a trailing standalone line.

    Example desired output:
      ... has been changed From: White PVC U:94V-0 To: White PVC FM 4910.
    """
    if not text:
        return text

    lines = [ln.strip() for ln in str(text).splitlines() if ln.strip()]
    if not lines:
        return ""

    out: List[str] = []
    i = 0
    while i < len(lines):
        line = lines[i]

        # Case 1: single-line "From: X To: Y"
        m_inline = re.match(
            r"^From\s*:\s*(.+?)\s+To\s*:\s*(.+?)\s*\.?$",
            line,
            flags=re.IGNORECASE,
        )
        if m_inline:
            from_val, to_val = _ensure_from_to_revision_order(m_inline.group(1).strip(), m_inline.group(2).strip())
            merged = f"From: {from_val} To: {to_val}"
            if out:
                out[-1] = out[-1].rstrip(" .") + f" {merged}."
            else:
                out.append(merged + ".")
            i += 1
            continue

        # Case 2: two-line pair "From: X" then "To: Y"
        m_from = re.match(r"^From\s*:\s*(.+?)\s*\.?$", line, flags=re.IGNORECASE)
        m_to = None
        if m_from and i + 1 < len(lines):
            m_to = re.match(r"^To\s*:\s*(.+?)\s*\.?$", lines[i + 1], flags=re.IGNORECASE)

        if m_from and m_to:
            from_val, to_val = _ensure_from_to_revision_order(m_from.group(1).strip(), m_to.group(1).strip())
            merged = f"From: {from_val} To: {to_val}"
            if out:
                out[-1] = out[-1].rstrip(" .") + f" {merged}."
            else:
                out.append(merged + ".")
            i += 2
            continue

        out.append(line)
        i += 1

    return "\n".join(out).strip()


def _remove_supplier_rev_tokens(text: str) -> str:
    """Remove supplier revision number mentions like Rev.02, Rev 01, rev 01."""
    if not text:
        return ""

    raw = str(text)
    # Remove tokens like Rev.02 / Rev 02 / rev 01 / REV-03 / REV: 04
    raw = re.sub(r"\bREV\.?\s*[-:]?\s*\d+\b", "", raw, flags=re.IGNORECASE)

    # Normalize spacing per-line to preserve section/newline formatting.
    lines = raw.splitlines()
    cleaned_lines: List[str] = []
    for ln in lines:
        ln = re.sub(r"[ \t]{2,}", " ", ln)
        ln = re.sub(r"\s+([,.;:])", r"\1", ln)
        ln = re.sub(r"([,.;:]){2,}", r"\1", ln)
        cleaned_lines.append(ln.rstrip())

    # Keep intentional blank lines between sections, collapse only excessive empties.
    out = "\n".join(cleaned_lines)
    out = re.sub(r"\n{3,}", "\n\n", out)
    return out.strip()


def _translate_sps_to_professional(text: str) -> str:
    """
    Converts raw supplier-voice SPS text to third-person professional language
    suitable for Applied Materials engineering documentation.
    Handles common first-person / informal phrases.
    """
    if not text:
        return text

    t = str(text)

    # First-person → third-person supplier references
    replacements = [
        # "we are unable" / "we cannot"
        (r"\bwe are unable to\b",        "Supplier is unable to"),
        (r"\bwe cannot\b",               "Supplier cannot"),
        (r"\bwe can not\b",              "Supplier cannot"),
        (r"\bwe are not able to\b",      "Supplier is unable to"),
        # "we are able" / "we can"
        (r"\bwe are able to\b",          "Supplier is able to"),
        (r"\bwe can\b",                  "Supplier can"),
        # "we have" / "we had" / "we will"
        (r"\bwe have\b",                 "Supplier has"),
        (r"\bwe had\b",                  "Supplier had"),
        (r"\bwe will\b",                 "Supplier will"),
        (r"\bwe would\b",                "Supplier would"),
        # "we are" general
        (r"\bwe are\b",                  "Supplier is"),
        # standalone "we" as subject
        (r"\bwe\b",                      "the Supplier"),
        # "our" possessive
        (r"\bour\b",                     "Supplier\'s"),
        # "us" object
        (r"\bus\b",                      "the Supplier"),
        # informal / vague phrases
        (r"\bplease\b",                  ""),
        (r"\bkindly\b",                  ""),
        (r"\bFYI\b",                     ""),
        (r"\bfor your information\b",    ""),
    ]

    for pattern, replacement in replacements:
        t = re.sub(pattern, replacement, t, flags=re.IGNORECASE)

    # Collapse any double-spaces introduced by blank replacements
    t = re.sub(r"  +", " ", t).strip(" ,;.")
    return t


def _to_single_line_summary(text: str, max_chars: int = 220) -> str:
    if not text:
        return ""

    normalized = re.sub(r"\s+", " ", str(text).strip())
    if not normalized:
        return ""

    sentence_parts = re.split(r"(?<=[.!?])\s+", normalized)
    single_line = sentence_parts[0].strip() if sentence_parts else normalized
    if len(single_line) > max_chars:
        single_line = single_line[: max_chars - 3].rstrip(" ,;:-") + "..."
    return single_line


def _build_reference_change_summary_from_payload(payload: Dict[str, Any]) -> str:
    pcr_lines: List[str] = []
    project_lines: List[str] = []
    esw_lines: List[str] = []

    for rec in payload.get("pcr_records") or []:
        pcr_id = rec.get("pcr_id")
        objective = _to_single_line_summary(rec.get("problem") or "")
        if pcr_id and objective:
            pcr_lines.append(f"PCR#   {pcr_id}: {objective}")

    for rec in payload.get("project_records") or []:
        project_id = rec.get("project_id")
        objective = _to_single_line_summary(rec.get("defined_scope") or rec.get("project_name") or "")
        if project_id and objective:
            project_lines.append(f"Project#   {project_id}: {objective}")

    for rec in payload.get("esw_records") or []:
        ec_number = rec.get("ec_number")
        objective = _to_single_line_summary(rec.get("title") or "")
        if ec_number and objective:
            esw_lines.append(f"ESW#   {ec_number}: {objective}")

    # SPS records are intentionally excluded from the Reference Change Summary.
    # SPS# appears once only — in the Problem Description opening sentence generated by the LLM.

    groups: List[str] = []
    if pcr_lines:
        groups.append("\n".join(_dedupe_lines(pcr_lines)))
    if project_lines:
        groups.append("\n".join(_dedupe_lines(project_lines)))
    if esw_lines:
        groups.append("\n".join(_dedupe_lines(esw_lines)))

    return "\n\n".join(groups).strip()


def _normalize_bullets(text: str, crisp: bool = False) -> str:
    if not text:
        return ""
    if crisp:
        return _short_crisp_bullets(text, max_bullets=3, max_words=12)

    # Normalize mixed numbering/bullet styles into clean bullet points.
    raw = str(text)
    raw = re.sub(r"\s+(?=\d+\.\s+)", "\n", raw)

    lines = []
    for raw_line in raw.splitlines():
        ln = raw_line.strip()
        ln = re.sub(r"^[-*•]\s*", "", ln)
        ln = re.sub(r"^\(?\d+\)?[.)]\s*", "", ln)
        ln = re.sub(r"^[A-Za-z][.)]\s*", "", ln)
        if re.fullmatch(r"\d+\.?", ln):
            continue
        ln = re.sub(r"\s+", " ", ln).strip(" -–—:;,")
        if ln:
            if not re.search(r"[.!?]$", ln):
                ln = ln + "."
            lines.append(f"- {ln}")
    return "\n".join(_dedupe_lines(lines)).strip()


def _extract_scope_part_descriptions(source_text: str) -> Dict[str, str]:
    """
    Parses 'Scope Parts:' block in source_text and returns a mapping of
    part_number -> part_description for entries like '0150-43377 - RING, FOCUS'.
    """
    pn_desc: Dict[str, str] = {}
    in_scope = False
    for raw in (source_text or "").splitlines():
        line = raw.strip()
        if re.match(r"^Scope\s+Parts\s*:", line, flags=re.IGNORECASE):
            in_scope = True
            continue
        if in_scope:
            if not line:
                continue
            # Stop if we hit another top-level section header
            if re.match(r"^[A-Za-z][A-Za-z\s]+:", line) and not line.startswith("-"):
                break
            entry = re.sub(r"^[-*•]\s*", "", line).strip()
            pns = extract_part_numbers(entry)
            if pns:
                pn = pns[0]
                desc = re.sub(rf"\b{re.escape(pn)}\b", "", entry, flags=re.IGNORECASE)
                desc = re.sub(r"\s+", " ", desc).strip(" -–—:;,.")
                if desc:
                    pn_desc[pn] = desc
    return pn_desc


def _extract_labeled_part_descriptions(source_text: str) -> Dict[str, str]:
    """Parse labeled part fields like 'Part Number:' and 'Part Description:' from input text."""
    pn_desc: Dict[str, str] = {}
    current_pn = ""

    for raw in (source_text or "").splitlines():
        line = raw.strip()
        if not line:
            continue

        m_pn = re.search(
            rf"\bPart\s*Number\s*:\s*({PN_FRAGMENT})\b",
            line,
            flags=re.IGNORECASE,
        )
        if m_pn:
            current_pn = m_pn.group(1).upper()

        m_desc = re.search(r"\bPart\s*Description\s*:\s*(.+)$", line, flags=re.IGNORECASE)
        if m_desc and current_pn:
            desc = re.sub(r"\s+", " ", m_desc.group(1)).strip(" -:;,.")
            if desc and current_pn not in pn_desc:
                pn_desc[current_pn] = desc

    return pn_desc


def _build_part_description_map(source_text: str) -> Dict[str, str]:
    """Combine known part-description sources for deterministic formatting."""
    out = _extract_scope_part_descriptions(source_text)
    labeled = _extract_labeled_part_descriptions(source_text)
    for pn, desc in labeled.items():
        if pn not in out and desc:
            out[pn] = desc
    return out


def _unique_preserve(items: List[str]) -> List[str]:
    seen = set()
    out: List[str] = []
    for item in items:
        val = str(item or "").strip().upper()
        if val and val not in seen:
            seen.add(val)
            out.append(val)
    return out


def _derive_part_sections_from_bom(
    normalized_affected: str,
    normalized_issue: str,
    source_text: str,
) -> Tuple[str, str]:
    """Ensure Issue part tracks child issue part; Affected part tracks parent part."""
    rels = extract_bom_relationships(source_text or "")
    if not rels:
        return normalized_affected, normalized_issue

    parent_to_children: Dict[str, List[str]] = {}
    child_to_parents: Dict[str, List[str]] = {}
    for parent, child in rels:
        parent_to_children.setdefault(parent, []).append(child)
        child_to_parents.setdefault(child, []).append(parent)

    affected_pns = _unique_preserve(extract_part_numbers(normalized_affected))
    issue_pns = _unique_preserve(extract_part_numbers(normalized_issue))

    # If Issue is blank, infer it from children of selected parents.
    if not issue_pns and affected_pns:
        inferred_issue: List[str] = []
        for parent in affected_pns:
            inferred_issue.extend(parent_to_children.get(parent, []))
        issue_pns = _unique_preserve(inferred_issue)

    # If Affected is blank, infer it from parents of selected issue parts.
    if not affected_pns and issue_pns:
        inferred_affected: List[str] = []
        for child in issue_pns:
            inferred_affected.extend(child_to_parents.get(child, []))
        affected_pns = _unique_preserve(inferred_affected)

    # If both are blank, derive complete mapping from BOM relationships.
    if not affected_pns and not issue_pns:
        affected_pns = _unique_preserve([p for p, _ in rels])
        issue_pns = _unique_preserve([c for _, c in rels])

    affected_text = ", ".join(affected_pns) if affected_pns else normalized_affected
    issue_text = ", ".join(issue_pns) if issue_pns else normalized_issue
    return affected_text, issue_text


def _dedupe_problem_sentences(text: str, max_lines: int = 0) -> str:
    """Remove near-duplicate lines from Problem Description while preserving concise detail."""
    if not text:
        return ""

    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    if not lines:
        return ""

    def _norm(s: str) -> str:
        return re.sub(r"\s+", " ", re.sub(r"[^a-z0-9 ]", " ", s.lower())).strip()

    def _tokens(s: str) -> set:
        return {t for t in _norm(s).split() if t}

    def _looks_duplicate(a: str, b: str) -> bool:
        na, nb = _norm(a), _norm(b)
        if not na or not nb:
            return False
        if na == nb or na in nb or nb in na:
            return True

        ta, tb = _tokens(a), _tokens(b)
        if not ta or not tb:
            return False
        overlap = len(ta & tb) / float(min(len(ta), len(tb)))
        return overlap >= 0.8

    kept: List[str] = []
    for ln in lines:
        if not kept:
            kept.append(ln)
            continue
        if any(_looks_duplicate(existing, ln) for existing in kept):
            continue
        kept.append(ln)
        if max_lines > 0 and len(kept) >= max_lines:
            break

    return "\n".join(kept).strip()


def _break_long_sentences(text: str, max_length: int = 180, max_words: int = 30) -> str:
    """Intelligently break overly long, complex sentences into 2-3 shorter, clearer sentences."""
    if not text:
        return ""

    text = str(text).strip()
    
    # If text is short enough, return as-is
    if len(text) <= max_length:
        return text
    
    sentences = re.split(r'(?<=[.!?])\s+', text)
    result_sentences: List[str] = []
    
    for sent in sentences:
        sent = sent.strip()
        if not sent:
            continue
        
        # If sentence is short enough, keep it as-is
        if len(sent) <= max_length:
            result_sentences.append(sent)
            continue
        
        # Try to break long sentence at semicolons or logical break points
        word_count = len(sent.split())
        if word_count <= max_words:
            result_sentences.append(sent)
            continue
        
        # Look for break points
        broken = False
        
        # Try splitting on semicolons
        if ';' in sent:
            parts = [p.strip() for p in sent.split(';')]
            result_sentences.extend([p + '.' if not p.endswith(('.', '!', '?')) else p for p in parts if p])
            broken = True
        
        # Try splitting before common phrases that start secondary clauses
        if not broken:
            break_patterns = [
                r'(\,\s*(?:resulting\s+in|causing|preventing|leading\s+to|thereby|which)\s+)',
                r'(\,\s*(?:however|but|yet|additionally|further|furthermore|moreover)\s+)',
                r'(\s+(?:Currently|Only|Specifically)\s+)',
            ]
            
            for pattern in break_patterns:
                if re.search(pattern, sent, flags=re.IGNORECASE):
                    match = re.search(pattern, sent, flags=re.IGNORECASE)
                    if match:
                        pos = match.start()
                        before = sent[:pos].rstrip(', ')
                        after = sent[pos:].lstrip(' ,')
                        if after.lower().startswith(('resulting', 'causing', 'preventing', 'leading', 'thereby', 'which', 'however', 'but', 'yet')):
                            # Capitalize first letter of second sentence
                            after = after[0].upper() + after[1:] if after else after
                        
                        # Ensure both parts end correctly
                        if before and not before.endswith(('.', '!', '?')):
                            before = before + '.'
                        if after and not after.endswith(('.', '!', '?')):
                            after = after + '.'
                        
                        result_sentences.append(before)
                        result_sentences.append(after)
                        broken = True
                        break
        
        # If no break point found, split after every ~25-30 words
        if not broken:
            words = sent.split()
            chunk_size = max_words - 5
            for i in range(0, len(words), chunk_size):
                chunk = ' '.join(words[i:i+chunk_size])
                if not chunk.endswith(('.', '!', '?')):
                    chunk += '.'
                result_sentences.append(chunk)
            broken = True
        
        if not broken:
            result_sentences.append(sent)
    
    return '\n'.join(result_sentences).strip()


def _inject_single_part_description(part_text: str, source_text: str) -> str:
    """If a part section has one unique part number and we know its description, append it."""
    if not part_text:
        return ""

    pns = extract_part_numbers(part_text)
    unique_pns = list(dict.fromkeys(pns))
    if len(unique_pns) != 1:
        return part_text

    pn = unique_pns[0]
    desc_map = _build_part_description_map(source_text)
    desc = (desc_map.get(pn) or "").strip()
    if not desc:
        return part_text

    return f"{pn} - {desc}"


def _extract_short_title_input(source_text: str) -> str:
    m = re.search(r"^\s*Short\s+Title\s+Input\s*:\s*(.+)$", source_text or "", flags=re.IGNORECASE | re.MULTILINE)
    return m.group(1).strip() if m else ""


def _remove_known_part_descriptions(text: str, source_text: str, remove_standalone: bool = False) -> str:
    """Strip known part descriptions from free-form text while preserving part numbers."""
    if not text:
        return ""

    cleaned = str(text)
    desc_map = _build_part_description_map(source_text)
    if not desc_map:
        return cleaned.strip()

    for pn, desc in sorted(desc_map.items(), key=lambda item: len(item[1]), reverse=True):
        desc = str(desc or "").strip()
        if not desc:
            continue

        pn_pat = re.escape(pn)
        desc_pat = re.escape(desc)

        cleaned = re.sub(
            rf"(\b(?:P/N|PN|Part\s*Number)\s*[:#-]?\s*{pn_pat}\b)\s*(?:[-–—,;:]|\()\s*{desc_pat}\s*(?:\))?",
            r"\1",
            cleaned,
            flags=re.IGNORECASE,
        )
        cleaned = re.sub(
            rf"(\b{pn_pat}\b)\s*(?:[-–—,;:]|\()\s*{desc_pat}\s*(?:\))?",
            r"\1",
            cleaned,
            flags=re.IGNORECASE,
        )

        if remove_standalone:
            cleaned = re.sub(
                rf"(?<![A-Za-z0-9]){desc_pat}(?![A-Za-z0-9])",
                "",
                cleaned,
                flags=re.IGNORECASE,
            )

    cleaned = re.sub(r"\s+", " ", cleaned)
    cleaned = re.sub(r"\s+([,.;:])", r"\1", cleaned)
    cleaned = re.sub(r"\(\s*\)", "", cleaned)
    cleaned = re.sub(r"\s*[-–—,:;]\s*[-–—,:;]\s*", " - ", cleaned)
    cleaned = re.sub(r"\s{2,}", " ", cleaned)
    return cleaned.strip(" -–—:;,.")


def _derive_complete_title_body(source_text: str) -> str:
    """Build a clearer fallback title body when the model returns a clipped fragment."""
    candidates: List[str] = []

    short_title = _extract_short_title_input(source_text)
    if short_title:
        candidates.append(short_title)

    patterns = [
        r"^\s*Title\s*:\s*(.+)$",
        r"^\s*Key\s+Objective\s*:\s*(.+)$",
        r"^\s*Problem\s+Statement\s+Input\s*:\s*(.+)$",
        r"^\s*Proposed\s+Solution\s+Input\s*:\s*(.+)$",
        r"^\s*Supplier-reported issue .*?:\s*(.+)$",
    ]
    for pattern in patterns:
        candidates.extend(re.findall(pattern, source_text or "", flags=re.IGNORECASE | re.MULTILINE))

    for raw in candidates:
        body = _remove_known_part_descriptions(raw, source_text, remove_standalone=True)
        body = re.sub(r"\bPart\s*Number\b", "", body, flags=re.IGNORECASE)
        body = re.sub(rf"\b{PN_FRAGMENT}\b", "", body, flags=re.IGNORECASE)
        body = re.sub(
            r"^\s*(?:Per\s+SPS#\s*\d+[, ]*)?(?:Supplier\s+notified(?:\s+Applied\s+Materials)?\s+that\s+)?",
            "",
            body,
            flags=re.IGNORECASE,
        )
        body = re.split(r"(?<=[.!?])\s+", body, maxsplit=1)[0]
        body = re.split(
            r"\b(?:resulting\s+in|causing|preventing|leading\s+to|thereby|because|since|so\s+that|currently|only)\b",
            body,
            maxsplit=1,
            flags=re.IGNORECASE,
        )[0]
        body = re.sub(r"\s+", " ", body).strip(" -–—:;,.")
        if body and not _title_looks_incomplete(body):
            return body

    return ""


def _trim_title_words(text: str, max_len: int = 75) -> str:
    t = re.sub(r"\s+", " ", str(text or "")).strip(" -–—:;,.")
    if len(t) <= max_len:
        return t

    words = t.split()
    out = ""
    for w in words:
        candidate = w if not out else f"{out} {w}"
        if len(candidate) > max_len:
            break
        out = candidate
    return out if out else t[:max_len].rstrip(" -–—:;,.")


def _title_looks_incomplete(text: str) -> bool:
    t = str(text or "").strip().lower()
    if not t:
        return True
    # Pure number or single/very-short token with no engineering meaning.
    if re.fullmatch(r"\d+", t):
        return True
    if len(t) < 5:
        return True
    if t.endswith(("-", "/", ":", ";", ",", "(")):
        return True
    if re.search(r"\b(?:and|or|with|for|to)\s*$", t):
        return True
    if re.search(r"\bnote\s+\d+\s*$", t):
        return True
    # Common clipped suffixes seen when text is cut mid-word.
    bad_suffixes = (
        "repla", "replac", "modif", "obsolescen", "specifi", "substitu", "materia", "procuremen",
    )
    return any(t.endswith(sfx) for sfx in bad_suffixes)


def _normalize_title_with_part_number(title: str, source_text: str) -> str:
    """
    Enforces title format: <PART_NUMBER> - <short technical title>
    and removes 'Part Number' wording and part descriptions from title.
    """
    t = str(title or "").strip()
    if not t:
        return ""

    title_pns = extract_part_numbers(t)
    source_pns = extract_part_numbers(source_text)
    # Never keep more than one part number in Title.
    lead_pn = source_pns[0] if source_pns else ""

    # Start by removing ALL PNs from body, then optionally prepend one clean lead PN.
    body = re.sub(rf"\bPart\s*Number\s*[:\-]?\s*{PN_FRAGMENT}\b", "", t, flags=re.IGNORECASE)
    body = re.sub(rf"\b{PN_FRAGMENT}\b", "", body, flags=re.IGNORECASE)
    body = re.sub(r"\bPart\s*Number\b", "", body, flags=re.IGNORECASE)
    body = re.sub(r"\s*[-–—]\s*[-–—]+\s*", " - ", body)
    body = re.sub(r"\s+", " ", body).strip(" -–—:;,.")
    body = _remove_known_part_descriptions(body, source_text, remove_standalone=True)

    # Remove part description (item name) from the title body if it appears there.
    pn_desc_map = _extract_scope_part_descriptions(source_text)
    part_desc = pn_desc_map.get(lead_pn, "")
    if part_desc:
        body = re.sub(
            rf"^{re.escape(part_desc)}\s*[-–—]?\s*",
            "",
            body,
            flags=re.IGNORECASE,
        ).strip(" -–—:;,.")

    # If body is empty, fallback to Short Title input from payload composition.
    if not body:
        body = _extract_short_title_input(source_text)
        body = re.sub(rf"\b{PN_FRAGMENT}\b", "", body, flags=re.IGNORECASE)
        body = re.sub(r"\s+", " ", body).strip(" -–—:;,.")
        body = _remove_known_part_descriptions(body, source_text, remove_standalone=True)

    if not body or _title_looks_incomplete(body):
        fallback_body = _derive_complete_title_body(source_text)
        if fallback_body:
            body = fallback_body

    if not body:
        body = "Engineering Change Update"

    # If the draft had multiple PNs or the sentence looks clipped, do not prefix PN.
    keep_pn_prefix = bool(lead_pn) and len(set(title_pns)) <= 1 and not _title_looks_incomplete(body)

    body_limit = max(20, 75 - len(lead_pn) - 3) if keep_pn_prefix else 75
    body = _trim_title_words(body, max_len=body_limit)
    candidate = f"{lead_pn} - {body}" if keep_pn_prefix else body

    # Guard against trimmed clipped-looking result.
    if _title_looks_incomplete(candidate):
        fallback_body = _derive_complete_title_body(source_text) or body
        no_pn = _trim_title_words(fallback_body, max_len=75)
        return no_pn or "Engineering Change Update"

    return candidate or "Engineering Change Update"


def _apply_engineering_shortcuts(text: str) -> str:
    """
    Replaces known engineering full terms with approved shortcuts.
    """
    if not text:
        return ""

    out = str(text)
    for full_term, shortcut in sorted(ENGINEERING_SHORTCUTS, key=lambda x: len(x[0]), reverse=True):
        pattern = rf"(?<!\w){re.escape(full_term)}(?!\w)"
        out = re.sub(pattern, shortcut, out, flags=re.IGNORECASE)
    return out


# ----------------------------
# Validation
# ----------------------------
def validate_engineering_output(output: str, source_text: str) -> List[str]:
    violations: List[str] = []

    if not re.search(r"^\s*Title\s*:", output, flags=re.IGNORECASE | re.MULTILINE):
        violations.append("Missing 'Title:' line.")
    if not re.search(r"^\s*Problem Statement\s*:", output, flags=re.IGNORECASE | re.MULTILINE):
        violations.append("Missing 'Problem Statement:' section.")
    if not re.search(r"^\s*(Solution Statement|Proposed Solution)\s*:", output, flags=re.IGNORECASE | re.MULTILINE):
        violations.append("Missing 'Solution Statement:' section.")

    title = _extract_title(output)
    if not title:
        violations.append("Title must not be empty.")

    if re.search(FIRST_PERSON_TERMS, output, flags=re.IGNORECASE):
        violations.append("Use third-person only.")

    for term in URGENCY_TERMS:
        if re.search(term, output, flags=re.IGNORECASE):
            violations.append("Do not use urgency language.")
            break

    for term in BLAME_TERMS:
        if re.search(term, output, flags=re.IGNORECASE):
            violations.append("Do not assign blame.")
            break

    problem = _extract_section(output, "Problem Statement")
    solution = _extract_section(output, "Solution Statement") or _extract_section(output, "Proposed Solution")

    if not problem:
        violations.append("Problem Statement section must not be empty.")
    if not solution:
        violations.append("Solution Statement section must not be empty.")

    if problem:
        affected = _extract_subsection(problem, "Affected/Impacted Part Numbers", PROBLEM_SUBSECTION_NAMES)
        if not affected:
            affected = _extract_subsection(problem, "Affected Part Numbers", PROBLEM_SUBSECTION_NAMES)
        issue_parts = _extract_subsection(problem, "Issue Part Number(s)", PROBLEM_SUBSECTION_NAMES) or \
                      _extract_subsection(problem, "Issue Part Numbers", PROBLEM_SUBSECTION_NAMES)
        problem_desc = _extract_subsection(problem, "Problem Description", PROBLEM_SUBSECTION_NAMES)
        impact = _extract_subsection(problem, "Impact Details", PROBLEM_SUBSECTION_NAMES)

        if not affected and not issue_parts:
            violations.append("Problem Statement must include part numbers in 'Affected/Impacted Part Numbers:' or 'Affected Part Numbers:'/'Issue Part Number(s):'.")
        if not problem_desc:
            violations.append("Problem Statement must include non-empty 'Problem Description:'.")
        if not impact:
            violations.append("Problem Statement must include non-empty 'Impact Details:'.")

        if problem_desc and not _has_bullet_lines(problem_desc):
            violations.append("Problem Description must use point-wise bullets with short complete sentences.")
        if impact and not _has_bullet_lines(impact):
            violations.append("Impact Details must use bullets.")

    if solution:
        solution_desc = _extract_subsection(solution, "Solution Description", SOLUTION_SUBSECTION_NAMES)
        if solution_desc and not _has_bullet_lines(solution_desc):
            violations.append("Solution Description must use point-wise bullets with short complete sentences.")

        if re.search(r"\bPart\s*Number\s*:", affected or "", flags=re.IGNORECASE):
            violations.append("Affected Part Numbers should not repeat 'Part Number:' labels.")
        if re.search(r"\bPart\s*Number\s*:", issue_parts or "", flags=re.IGNORECASE):
            violations.append("Issue Part Number(s) should not repeat 'Part Number:' labels.")

        for pat in DISALLOWED_PROBLEM_TERMS:
            if re.search(pat, "\n".join([problem_desc, impact]), flags=re.IGNORECASE):
                violations.append("Problem Statement contains solution language.")
                break

    source_parts = set(extract_part_numbers(source_text))
    out_parts = set(extract_part_numbers(output))
    if source_parts and not source_parts.issubset(out_parts):
        missing = sorted(source_parts - out_parts)
        violations.append(f"Missing part numbers from output: {', '.join(missing[:5])}")

    extra_parts = sorted(out_parts - source_parts)
    if extra_parts:
        violations.append(f"Output contains unknown part numbers: {', '.join(extra_parts[:5])}")

    return violations


# ----------------------------
# LLM Prompt + Call
# ----------------------------
def _call_databricks(prompt_text: str) -> str:
    if not DATABRICKS_URL or not DATABRICKS_API_KEY:
        return "Error: Databricks URL or API Key not set"

    headers = {
        "Authorization": f"Bearer {DATABRICKS_API_KEY}",
        "Content-Type": "application/json",
    }
    payload = {"messages": [{"role": "user", "content": prompt_text.strip()}]}

    try:
        resp = requests.post(DATABRICKS_URL, headers=headers, json=payload, timeout=60)
    except requests.RequestException as e:
        return f"Error: Request failed - {e}"

    if resp.status_code != 200:
        return f"Error: {resp.status_code} - {resp.text}"

    try:
        data = resp.json()
        content = data["choices"][0]["message"]["content"]
        if not content:
            return "Error: LLM returned empty content"
        return content
    except Exception:
        return f"Error: Unexpected response format - {resp.text}"


def _build_engineering_prompt(user_text: str) -> str:
    return f"""
You are acting as a Manufacturing / NPI Engineer.
Draft a formal, objective, audit-ready ECR response.

STRICT RULES:
- Use only input facts.
- Third-person only.
- No urgency, no blame, no opinions.
- Do not mix problem and solution content.
- Do not include revision notation (Rev, Rev., Rev 01, etc.) in Problem Statement or Solution Statement.
- SOURCE FIELD MAPPING — strictly observe which input fields feed which output section:
    • Problem Statement draws ONLY from: 'Problem:', 'Supplier-reported issue:', 'Defined Scope (Problem Basis):', 'Key Objective:' fields.
    • Solution Statement draws ONLY from: 'Solution:', 'Supporting Solution/Change:', 'Deliverables (Solution Basis):', 'Proposed Solution from Email:', 'User Proposed Solution:' fields.
    • NEVER copy a Solution/Proposed Change field value into Problem Statement.
    • NEVER copy a Problem/Issue field value into Solution Statement.
- Understand and use engineering relationships when available in input:
  • BOM structure (Parent/Child part relationship; Child reports to Parent).
  • Obsolete-to-replacement mapping (which part is being obsoleted and which replaces it).
  • Use these relationships to write clearer and accurate Problem/Solution sentences.
- If the input contains "Email Content:", read the full email body to extract the problem, root cause, and proposed solution.
- If the input also contains "User Entered Problem:" or "User Proposed Solution:", use those to supplement or refine what was found in the email.
- Ignore any alternative, rejected, or not-approved solutions mentioned in the email; use only the final approved/decided solution.

TITLE:
- Short and technical.
- If multiple affected part numbers exist, do not list multiple part numbers in title.
- If a valid Applied Materials part number is present in the source data, start title with it: <PART_NUMBER> - <short technical title>
- If NO part number is found in the source data, write the title as plain text WITHOUT any part number prefix.
- Do NOT invent, fabricate, or guess part numbers. Only use part numbers explicitly present in the input.
- Do not write the words "Part Number" in title.
- Do NOT include the part description (e.g. item name such as "RING, FOCUS" or "SCREW, M6") in the title. The title after the part number must describe the engineering issue or change, not the part name.
- Title must be clear, complete, and understandable within 75 characters.
- Never include more than one part number in title.
- If the title would look clipped, incomplete, or unclear with a part number prefix, remove part number(s) and output a clear plain-language title.

PROBLEM STATEMENT:
Must include exact subsection order:

Reference Change Summary:
- If reference records are available, this subsection must appear first.
- Use one line per reference record for PCR, Project, and ESW records only.
- Format each line exactly as identifier followed by colon and a concise objective/change summary.
- Preferred patterns:
    PCR <number>: <key objective of PCR>
    Project <number>: <key objective of Project>
    ESW <number>: <key objective of ESW>
- Do NOT write an SPS line in Reference Change Summary. The SPS reference belongs in Problem Description (see below).
- Keep each line concise and factual.
- IMPORTANT — SPS context: SPS is a supplier communication platform, NOT an Applied Materials engineering document.
  The ECR describes changes to AM engineering documents, not changes to SPS.
  SPS content is written in the supplier's voice (first-person, informal). You MUST translate it:
    • Convert all first-person ("we", "our", "us") to third-person ("Supplier", "Supplier's").
    • Replace informal phrases with professional equivalents
      (e.g. "we are able to purchase" → "Supplier is able to procure",
             "we cannot get" → "Supplier is unable to procure").
    • The Problem Description MUST open with:
      "Per SPS# <number>, Supplier notified that <translated professional summary of the issue>."
      Then continue with additional engineering detail as needed.
    • If SPS records contain Part Number and Part Description fields, include those part numbers in the Affected/Impacted Part Numbers or Issue Part Number(s) section as applicable.
  The SPS number must appear ONLY ONCE in the entire output — in this opening sentence of Problem Description.
  Do NOT write SPS# in Reference Change Summary, Solution Statement, or anywhere else in the narrative.

 Affected/Impacted Part Numbers (use this header if Affected and Issue are same):
 - List only Applied Materials part numbers found in the input.
 - AM part number format: 4-5 digits, hyphen, exactly 5 digits (e.g. NNNN-NNNNN or ESWNNN-NNNNN where N is a digit).
 - Do NOT invent, guess, or use example/placeholder part numbers. Only copy part numbers explicitly present in the source input.
 - If no AM part number is found in the input, write only: -
 - Do not include SPS IDs, PCR IDs, Project IDs, QN IDs, or supplier/manufacturer part numbers.
 - If multiple part numbers exist, list as comma-separated values on one line.
 - Do NOT repeat "Part Number:" label in this subsection.
 - CRITICAL LOGIC FOR PART DESCRIPTIONS:
     * If exactly ONE part number is identified in ALL input records combined:
         - AND that part number has a description shown in the input (from SPS records labeled "Part Number: X" and "Part Description: Y"):
         - THEN include the description: <PART_NUMBER> - <Part Description>
         - Example: 0042-90060 - RING, FOCUS
     * If TWO OR MORE part numbers exist anywhere in the input: List part numbers ONLY, no descriptions.
     * If ONE part number exists but NO description found: List part number ONLY, no description.

Affected Part Numbers (use this header only if different from Issue):
- List only Applied Materials part numbers found in the input. No descriptions.
- AM part number format: 4-5 digits, hyphen, exactly 5 digits (e.g. NNNN-NNNNN or ESWNNN-NNNNN where N is a digit).
- Do NOT invent, guess, or use example/placeholder part numbers. Only copy part numbers explicitly present in the source input.
- If no AM part number is found in the input, write only: -
- Do not include SPS IDs, PCR IDs, Project IDs, QN IDs, or supplier/manufacturer part numbers.
- If similar scope has multiple part numbers, list as comma-separated values on one line.
- Do NOT repeat "Part Number:" label in this subsection.

Issue Part Number(s) (use this header only if different from Affected):
- List only issue-causing Applied Materials part numbers found in the input.
- Issue Part Number(s) must be the child part number(s) that carry the issue.
- AM part number format: 4-5 digits, hyphen, exactly 5 digits (e.g. NNNN-NNNNN or ESWNNN-NNNNN where N is a digit).
- Do NOT invent, guess, or use example/placeholder part numbers. Only copy part numbers explicitly present in the source input.
- If no AM part number is found in the input, write only: -
- Do not include SPS IDs, PCR IDs, Project IDs, QN IDs, or supplier/manufacturer part numbers.
- If multiple similar-scope issue parts exist, list as comma-separated values.
- Do NOT repeat "Part Number:" label in this subsection.
- Do NOT repeat "Part Number:" label in this subsection.
- CRITICAL LOGIC FOR PART DESCRIPTIONS:
    * If exactly ONE part number is identified in ALL input records combined:
        - AND that part number has a description shown in the input (from SPS records labeled "Part Number: X" and "Part Description: Y"):
        - THEN include the description: <PART_NUMBER> - <Part Description>
        - Example: 0042-90060 - RING, FOCUS
    * If TWO OR MORE part numbers exist anywhere in the input: List part numbers ONLY, no descriptions.
    * If ONE part number exists but NO description found: List part number ONLY, no description.

Affected Part Number(s) and Issue Part Number(s) Relationship Rule:
- If BOM parent/child context exists, set Affected Part Number(s) to parent part number(s).
- If BOM parent/child context exists, set Issue Part Number(s) to child issue part number(s).
- Do not leave Issue Part Number(s) blank when an issue-carrying child part number is present.

Problem Description:
- Maximum 6 bullet points.
- Each bullet must be a single, complete sentence written entirely on one line. Never continue or split a sentence onto the next line.
- DESCRIBE WHAT IS WRONG: state the engineering issue, discrepancy, missing specification, or incorrect value factually.
- Use factual/descriptive language (e.g., "does not specify", "is incorrect", "is missing", "is undefined", "is not documented").
- Do NOT use action/solution verbs such as Update, Specify, Define, Add note, Revise, Clarify — those belong in Solution Description only.
- Do not include proposed fixes, corrective actions, or solution content.
- Do not include background, history, or narration.
- Do not repeat the part number, DWG reference (e.g. "Assy DWG for XXXX"), SPS number, or supplier narration across multiple bullets. Mention each at most once, or omit if already listed in part number sections above.
- Combine related issues into a single bullet where possible.
- No paragraphs; bullet points only.
- Output must be suitable for direct insertion into an ECN system.
- Do NOT include any "Reason Code" text (e.g. "Reason Code: Beyond Spec Request") in Problem Description.
- Do not mention part descriptions or item names; use part numbers only.
- If BOM Parent/Child is present, state the relationship concisely.
- If obsolete/replacement mapping is present, state which part is obsolete and which replaces it.
- If the output exceeds 6 bullets or repeats information, rewrite and compress it.

Impact Details:
- 2 to 3 short, crisp bullets only.
- Max ~12 words per bullet.

FROM / TO CHANGE FORMAT:
- In Problem Description, keep the From/To change inside the sentence (not as a separate trailing line).
- Preferred sentence pattern:
    Supplier notified that <what changed> From: <old value> To: <new value>.
- Do not place a standalone line like "From: ... To: ..." at the end of Problem Description.
- In Solution Description, From/To may be shown either inline or on separate lines.

SOLUTION STATEMENT:
Must include exact subsection order:

Solution Description:
- Maximum 6 bullet points.
- Each bullet must be a single, complete sentence written entirely on one line. Never continue or split a sentence onto the next line.
- DESCRIBE WHAT TO DO: state the corrective actions required to fix the problem.
- Use clear, direct, action-oriented imperative language (e.g., Update, Specify, Define, Add note, Revise, Clarify).
- Focus only on the actions required to resolve the issue — drawing updates, BOM changes, new notes, etc.
- Do NOT include problem descriptions, issue statements, or root-cause language — those belong in Problem Description only.
- Do not include background, history, or narration.
- Do not repeat the part number, DWG reference (e.g. "Assy DWG for XXXX"), SPS number, or supplier narration across multiple bullets. Mention each at most once.
- Combine related actions into a single bullet where possible.
- No paragraphs; bullet points only.
- Output must be suitable for direct insertion into an ECN system.
- Do not mention part descriptions or item names; use part numbers only.
- If BOM Parent/Child is present, describe the required BOM/document change concisely.
- If obsolete/replacement mapping is present, state the replacement action using that mapping.
- If "Proposed Solution from Email:" or "User Proposed Solution:" is provided, refine and integrate it here.
- If the output exceeds 6 bullets or repeats information, rewrite and compress it.

Benefits of the Proposed Solution:
- 2 to 3 short, crisp bullets only.
- Max ~12 words per bullet.

ENGINEERING SHORTCUT KEY USAGE:
- Use these shortcut keys whenever the corresponding full terms appear in Title, Problem Statement, or Solution Statement:
    Drawing = DWG
    Bill of Material = BOM
    Engineering Change Request = ECR
    Engineering Change Order = ECO
    Engineering Change Notice = ECN
    Revision = REV
    Fit, Form & Function = FFF
    Failure Mode & Effects Analysis = FMEA
    Root Cause Analysis = RCA
    Standard Operating Procedure = SOP
    First Article Inspection = FAI
    Lead Time = LT
    Purchase Order = PO
    Purchase Requisition = PR
    Approved Vendor List = AVL
    Minimum Order Quantity = MOQ
    Material Requirement Planning = MRP
    Goods Receipt = GR
    First In First Out = FIFO
    Manufacturing Drawing = Mfg. DWG
    Design FMEA = DFMEA
    Process FMEA = PFMEA
    Assembly = Assy
    Part Number = P/N
    Engineering Drawing = ENG DWG
    REV = Rev
    Obsolescence = Obs

OUTPUT FORMAT (exact top-level headers only):
Title:

Problem Statement:
Reference Change Summary:      (include when reference records exist)
Affected/Impacted Part Numbers:  (use when Affected and Issue are same)
Affected Part Number(s):            (use when different)
Issue Part Number(s):             (use when different)
Problem Description:
Impact Details:

Solution Statement:
Solution Description:
Benefits of the Proposed Solution:

INPUT:
<<<
{user_text}
>>>
""".strip()


def reframe_problem(user_text: str) -> str:
    prompt = _build_engineering_prompt(user_text)
    draft = _call_databricks(prompt)
    if draft.startswith("Error:"):
        return draft

    draft = clean_ai_output(draft)
    draft = remove_unwanted_sections(draft)
    draft = _remove_supplier_rev_tokens(draft)
    draft = re.sub(r"\bSPS\s*#?\s*\d+\b", "", draft, flags=re.IGNORECASE)

    def _critical_violations(items: List[str]) -> List[str]:
        return [
            v for v in items
            if any(kw in v for kw in ("SPS", "Missing", "must not be empty"))
        ]

    violations = validate_engineering_output(draft, user_text)
    critical_violations = _critical_violations(violations)
    if not critical_violations:
        return draft

    # Lightweight deterministic normalization pass (no extra LLM call).
    title = _normalize_title_with_part_number(_extract_title(draft).strip(), user_text)
    problem_section = _extract_section(draft, "Problem Statement").strip()
    solution_section = _extract_section(draft, "Solution Statement").strip() or _extract_section(draft, "Proposed Solution").strip()

    normalized_problem = _normalize_problem_statement(problem_section, user_text)
    normalized_solution = _normalize_solution_statement(solution_section)
    deterministic = (
        f"Title:\n{title}\n\n"
        f"Problem Statement:\n{normalized_problem}\n\n"
        f"Solution Statement:\n{normalized_solution}"
    ).strip()
    deterministic = clean_ai_output(deterministic)
    deterministic = remove_unwanted_sections(deterministic)
    deterministic = _remove_supplier_rev_tokens(deterministic)
    deterministic = re.sub(r"\bSPS\s*#?\s*\d+\b", "", deterministic, flags=re.IGNORECASE)

    deterministic_violations = validate_engineering_output(deterministic, user_text)
    deterministic_critical = _critical_violations(deterministic_violations)
    if len(deterministic_critical) < len(critical_violations):
        draft = deterministic
        critical_violations = deterministic_critical

    # Short corrective prompt fallback only when critical violations remain.
    if critical_violations:
        fix_prompt = f"""
Fix only these violations in the current draft:

- {"\n- ".join(critical_violations)}

Return the full corrected output using only these top-level headers:
Title:
Problem Statement:
Solution Statement:

Do not add explanations. Preserve existing valid content and edit minimally.

Current draft:
{draft}
""".strip()
        revised = _call_databricks(fix_prompt)
        if not revised.startswith("Error:"):
            revised = clean_ai_output(revised)
            revised = remove_unwanted_sections(revised)
            revised = _remove_supplier_rev_tokens(revised)
            revised = re.sub(r"\bSPS\s*#?\s*\d+\b", "", revised, flags=re.IGNORECASE)
            revised_violations = validate_engineering_output(revised, user_text)
            revised_critical = _critical_violations(revised_violations)
            if len(revised_critical) <= len(critical_violations):
                draft = revised

    return draft


# ----------------------------
# UI Payload Adapters
# ----------------------------
def _has_yes_reference(payload: Dict[str, Any]) -> bool:
    refs = payload.get("reference_inputs") or {}
    for _, v in refs.items():
        ans = str((v or {}).get("answer", "")).strip().lower()
        if ans == "yes":
            return True
    return False


def _any_tab_checked(payload: Dict[str, Any]) -> bool:
    flags = payload.get("include_tabs_flags") or {}
    return any(bool(v) for v in flags.values())


def _compose_user_text(payload: Dict[str, Any]) -> str:
    chunks: List[str] = []

    short_title = (payload.get("short_title") or "").strip()
    reason_code = (payload.get("reason_code") or "").strip()
    scope_parts = payload.get("scope_parts") or []
    refs = payload.get("reference_inputs") or {}
    current_problem = (payload.get("current_problem_text") or "").strip()
    current_solution = str(
        payload.get("current_solution_text") or payload.get("proposed_solution_text") or ""
    ).strip()
    selected_tab_payload = payload.get("selected_tab_payload") or {}

    # ------------------------------------------------------------------
    # PCR-sourced context (Question 1 – Databricks-enriched path)
    # ------------------------------------------------------------------
    pcr_records = payload.get("pcr_records") or []
    derived_psn = payload.get("derived_psn") or {}
    skipped_pcrs = payload.get("skipped_pcrs") or []
    project_records = payload.get("project_records") or []
    skipped_projects = payload.get("skipped_projects") or []
    sps_records = payload.get("sps_records") or []
    esw_records = payload.get("esw_records") or []
    reference_ecr_records = payload.get("reference_ecr_records") or []
    qn_records = payload.get("qn_records") or []

    if pcr_records:
        pcr_lines: List[str] = []
        for rec in pcr_records:
            pcr_id   = rec.get("pcr_id", "")
            status   = rec.get("status", "")
            problem  = (rec.get("problem") or "").strip()
            solution = (rec.get("solution") or "").strip()
            psn      = rec.get("psnnumber") or ""
            entry = (
                f"PCR {pcr_id} (Status: {status}):\n"
                f"  Problem: {problem}\n"
                f"  Solution: {solution}"
            )
            if psn:
                entry += f"\n  Associated PSN: {psn}"
            pcr_lines.append(entry)
        chunks.append("PCR Records from Databricks:\n\n" + "\n\n".join(pcr_lines))

    if derived_psn.get("answer") == "Yes":
        psn_nums = ", ".join(str(n) for n in (derived_psn.get("numbers") or []))
        chunks.append(f"PSN: Yes – PSN Number(s): {psn_nums}")
    elif derived_psn.get("answer") == "No":
        chunks.append("PSN: No – No associated PSN found in the PCR records.")

    if skipped_pcrs:
        skip_lines = [
            f"  PCR {s.get('pcr_id')} – Status: {s.get('status')}"
            for s in skipped_pcrs
        ]
        chunks.append(
            "Note – the following PCR(s) were skipped (inactive status) "
            "and are not included in the summary:\n" + "\n".join(skip_lines)
        )

    if project_records:
        project_lines: List[str] = []
        for rec in project_records:
            project_id = rec.get("project_id", "")
            status = rec.get("status", "")
            name = (rec.get("project_name") or "").strip()
            defined_scope = (rec.get("defined_scope") or "").strip()
            deliverables = (rec.get("deliverables") or "").strip()
            entry = (
                f"Project {project_id} (Status: {status}):\n"
                f"  Name: {name}\n"
                f"  Defined Scope (Problem Basis): {defined_scope}\n"
                f"  Deliverables (Solution Basis): {deliverables}"
            )
            project_lines.append(entry)
        chunks.append("Project Records from Databricks:\n\n" + "\n\n".join(project_lines))

    if skipped_projects:
        skip_project_lines = [
            f"  Project {s.get('project_id')} - Status: {s.get('status')}"
            for s in skipped_projects
        ]
        chunks.append(
            "Note - the following Project(s) were skipped (inactive status) "
            "and are not included in the summary:\n" + "\n".join(skip_project_lines)
        )

    def _append_ec_record_chunk(header: str, records: List[Dict[str, Any]], summary_label: str) -> None:
        if not records:
            return
        lines: List[str] = []
        for rec in records:
            ec_number = rec.get("ec_number") or rec.get("sps_id") or ""
            status = rec.get("status", "")
            title = (rec.get("title") or "").strip()
            problem = (rec.get("problem") or "").strip()
            solution = (rec.get("solution") or rec.get("proposed_solution") or "").strip()
            entry_lines = [f"{header} {ec_number} (Status: {status}):"]
            if title:
                entry_lines.append(f"  Title: {title}")
            entry_lines.append(f"  {summary_label}: {problem}")
            entry_lines.append(f"  Supporting Solution/Change: {solution}")
            lines.append("\n".join(entry_lines))
        chunks.append(f"{header} Records from Databricks:\n\n" + "\n\n".join(lines))

    def _append_sps_record_chunk(records: List[Dict[str, Any]]) -> None:
        """Like _append_ec_record_chunk but flags SPS as supplier-to-AM communication."""
        if not records:
            return
        lines: List[str] = []
        for rec in records:
            sps_id = rec.get("sps_id") or rec.get("ec_number") or ""
            status = rec.get("status", "")
            part_number = (rec.get("part_number") or "").strip()
            part_description = (rec.get("part_description") or "").strip()
            problem = (rec.get("problem") or "").strip()
            solution = (rec.get("solution") or rec.get("proposed_solution") or "").strip()
            entry = (
                f"SPS {sps_id} (Status: {status}) — SUPPLIER COMMUNICATION TO APPLIED MATERIALS:\n"
                f"  Supplier-reported issue (translate to third-person professional): {problem}\n"
                f"  Supplier-proposed change (translate to third-person professional): {solution}"
            )
            if part_number:
                entry += f"\n  Part Number: {part_number}"
            if part_description:
                entry += f"\n  Part Description: {part_description}"
            lines.append(entry)
        chunks.append(
            "SPS Records from Databricks (raw supplier voice — treat as supplier-to-AM communication):\n\n"
            + "\n\n".join(lines)
        )

    _append_sps_record_chunk(sps_records)
    _append_ec_record_chunk("ESW", esw_records, "Key Objective")
    _append_ec_record_chunk("Reference ECR", reference_ecr_records, "Key Objective")
    _append_ec_record_chunk("QN", qn_records, "Key Objective")
    # ------------------------------------------------------------------

    if short_title:
        chunks.append(f"Short Title Input: {short_title}")
    if reason_code:
        chunks.append(f"Reason Code: {reason_code}")

    if scope_parts:
        chunks.append("Scope Parts:\n" + "\n".join([f"- {str(x).strip()}" for x in scope_parts if str(x).strip()]))

    ref_lines = []
    for k, v in refs.items():
        ans = str((v or {}).get("answer", "")).strip()
        txt = str((v or {}).get("text", "")).strip()
        if ans or txt:
            ref_lines.append(f"{k} | Answer: {ans or 'N/A'} | Text: {txt}")
    if ref_lines:
        chunks.append("Reference Inputs:\n" + "\n".join(ref_lines))

    if current_problem:
        chunks.append("Problem Statement Input:\n" + current_problem)

    if current_solution:
        chunks.append("Proposed Solution Input:\n" + current_solution)

    if selected_tab_payload:
        chunks.append("Selected Tab Payload:\n" + json.dumps(selected_tab_payload, ensure_ascii=False, indent=2))

    text = "\n\n".join(chunks).strip()

    bom = build_bom_context(text)
    replacement = build_replacement_context(text)

    if bom:
        text = f"{text}\n\n{bom}"
    if replacement:
        text = f"{text}\n\n{replacement}"

    return text


def _rewrite_bullets_to_imperative(problem_desc: str, source_text: str = "") -> str:
    """
    Rewrites Problem Description bullet points from passive/negative phrasing
    into active imperative drawing-update statements, prefixed by a summary
    header.  Falls back to the original problem_desc on any LLM error.
    """
    if not problem_desc or not problem_desc.strip():
        return problem_desc

    # Detect primary part number from bullets first, then source context.
    pns = extract_part_numbers(problem_desc) or extract_part_numbers(source_text)
    primary_pn = pns[0] if pns else ""

    # Detect whether an SPS reference is present in the description.
    has_sps = bool(re.search(r"\bSPS#?\s*\d+", problem_desc, re.IGNORECASE))

    dwg_ref = f"Assy DWG for {primary_pn}" if primary_pn else "Assy DWG"

    prompt = (
        "You are an NPI/Manufacturing Engineer writing a formal Engineering Change Notice (ECN).\n\n"
        "TASK: Rewrite the bullet-point Problem Description below into action-oriented imperative "
        "statements suitable for a drawing update summary.\n\n"
        "RULES:\n"
        "1. One line per action. No bullet dashes, numbers, or list prefixes.\n"
        "2. Start each line with an imperative verb: Update, Specify, Define, Add note, "
        "Add handling note, Revise, Clarify, etc.\n"
        "3. Convert passive/negative sentences (\"X is not specified\") to active imperatives "
        "(\"Specify X as Y\").\n"
        "4. Keep all values, tolerances, part numbers, and technical specifics verbatim.\n"
        f"5. Do NOT include '{dwg_ref}' or any similar DWG reference in any output line — "
        "the header already identifies the drawing.\n"
        "6. Do NOT mention any person or human name in the output.\n"
        "7. If any text in the bullets is written in ALL CAPS (e.g. 'DO NOT BEND THE CORRUGATED TUBE'), "
        "preserve it exactly as-is — do not rephrase, lowercase, or reformat it.\n"
        "8. If the source context contains reference configurations or approval "
        "details relevant to a bullet, include them concisely on the same line.\n"
        "9. Output ONLY the imperative lines — no header, no preamble, no explanation, "
        "no blank lines between items.\n\n"
        "PROBLEM DESCRIPTION BULLETS:\n"
        f"{problem_desc.strip()}\n\n"
        "SOURCE CONTEXT (for enrichment only):\n"
        f"{(source_text or '').strip()[:2000]}"
    )

    result = _call_databricks(prompt)
    if not result or result.startswith("Error:"):
        return problem_desc  # graceful fallback

    # Strip any stray bullet/number prefixes from LLM output lines.
    lines: List[str] = []
    for raw in result.strip().splitlines():
        ln = raw.strip()
        ln = re.sub(r"^[-*\u2022]\s*", "", ln)
        ln = re.sub(r"^\d+[.)]\s*", "", ln)
        ln = ln.strip()
        if ln:
            lines.append(ln)

    if not lines:
        return problem_desc  # fallback when LLM returned nothing useful

    return "\n".join(lines)


def _normalize_problem_statement(problem_section: str, source_text: str = "") -> str:
    if not problem_section:
        return ""

    affected = _extract_subsection(problem_section, "Affected/Impacted Part Numbers", PROBLEM_SUBSECTION_NAMES)
    if not affected:
        affected = _extract_subsection(problem_section, "Affected Part Numbers", PROBLEM_SUBSECTION_NAMES)
    issue_parts = _extract_subsection(problem_section, "Issue Part Number(s)", PROBLEM_SUBSECTION_NAMES)
    if not issue_parts:
        issue_parts = _extract_subsection(problem_section, "Issue Part Numbers", PROBLEM_SUBSECTION_NAMES)
    problem_desc = _extract_subsection(problem_section, "Problem Description", PROBLEM_SUBSECTION_NAMES)
    impact_details = _extract_subsection(problem_section, "Impact Details", PROBLEM_SUBSECTION_NAMES)

    problem_desc = _compress_repeated_part_number_labels(problem_desc)
    # Never carry Reason Code metadata into Problem Description narrative.
    problem_desc = re.sub(
        r"(?im)^\s*[-*•]?\s*Reason\s*Code\s*:\s*.*$",
        "",
        problem_desc,
    )
    problem_desc = re.sub(
        r"(?i)\bReason\s*Code\s*:\s*Beyond\s*Spec\s*Request\b[\s,;:.\-]*",
        "",
        problem_desc,
    )
    problem_desc = _remove_known_part_descriptions(problem_desc, source_text, remove_standalone=True)
    problem_desc = _break_long_sentences(problem_desc, max_length=180, max_words=30)
    problem_desc = _linebreak_after_fullstop(problem_desc)
    problem_desc = _inline_problem_description_from_to(problem_desc)
    problem_desc = _dedupe_problem_sentences(problem_desc, max_lines=0)
    problem_desc = _normalize_bullets(problem_desc)

    normalized_affected = _normalize_multiline_list(affected) if affected else ""
    normalized_issue = _normalize_multiline_list(issue_parts) if issue_parts else ""

    normalized_affected, normalized_issue = _derive_part_sections_from_bom(
        normalized_affected,
        normalized_issue,
        source_text,
    )

    normalized_affected = _inject_single_part_description(normalized_affected, source_text)
    normalized_issue = _inject_single_part_description(normalized_issue, source_text)

    affected_pns = extract_part_numbers(normalized_affected)
    issue_pns = extract_part_numbers(normalized_issue)
    same_parts = bool(affected_pns and issue_pns and set(affected_pns) == set(issue_pns))

    blocks: List[str] = []
    if same_parts:
        blocks.append("Affected/Impacted Part Numbers:\n" + (normalized_affected or "-"))
    else:
        blocks.append("Affected Part Number(s):\n" + (normalized_affected or "-"))
        blocks.append("Issue Part Number(s):\n" + (normalized_issue or "-"))
    if problem_desc:
        blocks.append("Problem Description:\n" + problem_desc.strip())
    if impact_details:
        impact_details = _compress_repeated_part_number_labels(impact_details)
        blocks.append("Impact Details:\n" + _normalize_bullets(impact_details, crisp=True))

    return "\n\n".join(blocks).strip()


def _normalize_solution_statement(solution_section: str) -> str:
    if not solution_section:
        return ""

    proposed_change = _extract_subsection(solution_section, "Proposed Change", SOLUTION_SUBSECTION_NAMES)
    solution_desc = _extract_subsection(solution_section, "Solution Description", SOLUTION_SUBSECTION_NAMES)
    benefits = _extract_subsection(solution_section, "Benefits of the Proposed Solution", SOLUTION_SUBSECTION_NAMES)

    proposed_change = _compress_repeated_part_number_labels(proposed_change)
    proposed_change = _format_from_to_changes(proposed_change)
    solution_desc = _compress_repeated_part_number_labels(solution_desc)
    solution_desc = _remove_known_part_descriptions(solution_desc, solution_section, remove_standalone=True)
    solution_desc = _break_long_sentences(solution_desc, max_length=180, max_words=30)
    solution_desc = _linebreak_after_fullstop(solution_desc)
    solution_desc = _format_from_to_changes(solution_desc)
    solution_desc = _normalize_bullets(solution_desc)
    benefits = _compress_repeated_part_number_labels(benefits)

    # Backward-compatible parsing: if model still emits "Proposed Change",
    # fold it into a single "Solution Description" block.
    merged_solution_desc = "\n".join(
        [x.strip() for x in [proposed_change, solution_desc] if x and x.strip()]
    ).strip()
    merged_solution_desc = _normalize_bullets(merged_solution_desc)
    merged_solution_desc = _rewrite_bullets_to_imperative(merged_solution_desc)

    blocks: List[str] = []
    if merged_solution_desc:
        blocks.append("Solution Description:\n" + merged_solution_desc)
    if benefits:
        blocks.append("Benefits of the Proposed Solution:\n" + _normalize_bullets(benefits, crisp=True))

    return "\n\n".join(blocks).strip()


# ----------------------------
# Public Entry Points
# ----------------------------
def generate_full_pss(payload: Dict[str, Any]) -> Dict[str, str]:
    source_text = _compose_user_text(payload)
    # Apply source-based revision ordering early, before other normalization
    source_text = _apply_source_based_from_to_order(source_text, payload)
    if not source_text:
        return {"title": "", "problem_statement": "", "solution_statement": "", "raw": ""}

    drafted = reframe_problem(source_text)
    drafted = str(drafted or "").strip()

    if not drafted or drafted.startswith("Error:"):
        return {
            "title": str(payload.get("short_title") or "").strip()[:200],
            "problem_statement": (payload.get("current_problem_text") or source_text)[:4000],
            "solution_statement": str(payload.get("current_solution_text") or payload.get("proposed_solution_text") or "")[:4000],
            "raw": drafted or "",
        }

    title = _normalize_title_with_part_number(_extract_title(drafted).strip(), source_text)
    problem_section = _extract_section(drafted, "Problem Statement").strip()
    solution_section = _extract_section(drafted, "Solution Statement").strip()
    if not solution_section:
        solution_section = _extract_section(drafted, "Proposed Solution").strip()

    normalized_problem = _normalize_problem_statement(problem_section, source_text)
    normalized_solution = _normalize_solution_statement(solution_section)
    payload_reference_summary = _build_reference_change_summary_from_payload(payload)

    if payload_reference_summary:
        normalized_problem = (
            "Reference Change Summary:\n"
            + payload_reference_summary
            + ("\n\n" + normalized_problem if normalized_problem else "")
        ).strip()

    title = _apply_engineering_shortcuts(title)
    normalized_problem = _apply_engineering_shortcuts(normalized_problem)
    normalized_solution = _apply_engineering_shortcuts(normalized_solution)

    # Apply source-based From/To ordering correction before removing revision tokens
    normalized_problem = _apply_source_based_from_to_order(normalized_problem, payload)
    normalized_solution = _apply_source_based_from_to_order(normalized_solution, payload)
    
    # Remove supplier revision notation from Problem and Solution descriptions per requirements.
    normalized_problem = _remove_supplier_rev_tokens(normalized_problem)
    normalized_solution = _remove_supplier_rev_tokens(normalized_solution)

    # Fallback: if parsing extracted nothing despite the LLM returning content,
    # use the raw drafted text so the caller never gets all-empty fields.
    if not any([title, normalized_problem, normalized_solution]) and drafted:
        return {
            "title": str(payload.get("short_title") or "").strip()[:200],
            "problem_statement": drafted[:4000],
            "solution_statement": str(payload.get("current_solution_text") or payload.get("proposed_solution_text") or "")[:4000],
            "raw": drafted[:12000],
        }

    return {
        "title": title[:75],
        "problem_statement": normalized_problem[:4000],
        "solution_statement": normalized_solution[:4000],
        "raw": drafted[:12000],
    }


def generate_problem_summary(payload: Dict[str, Any]) -> str:
    return generate_full_pss(payload).get("problem_statement", "")[:4000]


def generate_pss(payload: Dict[str, Any]) -> str:
    return generate_problem_summary(payload)


def run_problem_summary(payload: Dict[str, Any]) -> str:
    return generate_problem_summary(payload)


# ----------------------------
# Email Parsing & Integration
# ----------------------------
def read_email_file(file_path: str) -> str:
    """
    Read reference attachment content and extract readable text.
    - .msg  : parsed via extract_msg (Outlook binary format)
    - .eml  : parsed via Python's email stdlib
    - .txt  : read as plain text
    - .pdf  : parsed via pypdf / PyPDF2
    - .pptx : parsed via python-pptx
    - .docx : parsed via python-docx
    - .ppt / .doc : best-effort plain-text fallback
    """
    if not file_path:
        return ""

    ext = os.path.splitext(file_path)[1].lower()

    # --- Outlook .msg ---
    if ext == ".msg":
        try:
            import extract_msg
            with extract_msg.openMsg(file_path) as msg:
                parts = []
                if msg.subject:
                    parts.append(f"Subject: {msg.subject.strip()}")
                if msg.sender:
                    parts.append(f"From: {msg.sender.strip()}")
                if msg.date:
                    parts.append(f"Date: {msg.date}")
                body = (msg.body or "").strip()
                if not body:
                    # fallback to HTML body stripped of tags
                    html = (msg.htmlBody or b"").decode("utf-8", errors="ignore")
                    body = re.sub(r"<[^>]+>", " ", html)
                    body = re.sub(r"[ \t]+", " ", body).strip()
                if body:
                    parts.append(f"\n{body}")
                return "\n".join(parts).strip()
        except Exception as e:
            return f"Error reading .msg file: {e}"

    # --- Standard .eml ---
    if ext == ".eml":
        try:
            import email as _email
            import email.policy
            with open(file_path, "rb") as f:
                msg = _email.message_from_binary_file(f, policy=_email.policy.default)
            parts = []
            for hdr in ("Subject", "From", "Date"):
                val = msg.get(hdr, "").strip()
                if val:
                    parts.append(f"{hdr}: {val}")
            # Walk MIME parts for plain text
            body_parts = []
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        charset = part.get_content_charset() or "utf-8"
                        body_parts.append(payload.decode(charset, errors="ignore"))
            if body_parts:
                parts.append("\n" + "\n".join(body_parts).strip())
            return "\n".join(parts).strip()
        except Exception as e:
            return f"Error reading .eml file: {e}"

    # --- PDF ---
    if ext == ".pdf":
        try:
            text_chunks = []
            try:
                from pypdf import PdfReader  # type: ignore
            except Exception:
                from PyPDF2 import PdfReader  # type: ignore

            reader = PdfReader(file_path)
            for page in reader.pages:
                txt = page.extract_text() or ""
                if txt.strip():
                    text_chunks.append(txt)
            return "\n\n".join(text_chunks).strip()
        except Exception as e:
            return f"Error reading .pdf file: {e}"

    # --- PowerPoint ---
    if ext == ".pptx":
        try:
            from pptx import Presentation  # type: ignore

            prs = Presentation(file_path)
            slide_text = []
            for idx, slide in enumerate(prs.slides, start=1):
                lines = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        t = str(shape.text).strip()
                        if t:
                            lines.append(t)
                if lines:
                    slide_text.append(f"Slide {idx}:\n" + "\n".join(lines))
            return "\n\n".join(slide_text).strip()
        except Exception as e:
            return f"Error reading .pptx file: {e}"

    # --- Word ---
    if ext == ".docx":
        try:
            from docx import Document  # type: ignore

            doc = Document(file_path)
            lines = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
            return "\n".join(lines).strip()
        except Exception as e:
            return f"Error reading .docx file: {e}"

    # --- Legacy office fallback ---
    if ext in {".doc", ".ppt"}:
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                txt = f.read().strip()
            if txt:
                return txt
            return (
                f"Error reading {ext} file: format is not directly supported. "
                f"Please save as .docx/.pptx and retry."
            )
        except Exception as e:
            return f"Error reading {ext} file: {e}"

    # --- Plain text / fallback ---
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read().strip()
    except Exception as e:
        return f"Error reading email file: {e}"


def extract_problem_from_email(email_content: str) -> str:
    """
    Parse email content and extract the problem statement.
    Looks for issue/problem/complaint descriptions in the email body.
    """
    if not email_content:
        return ""
    
    lines = email_content.split('\n')
    problem_lines = []
    in_body = False
    
    for i, line in enumerate(lines):
        # Skip email headers
        if line.strip() == "":
            in_body = True
            continue
        
        if not in_body:
            continue
        
        # Stop at signature or forwarded message
        if re.search(r'^(--|_____|==|From:|Sent:|To:|Cc:)', line):
            break
        
        problem_lines.append(line)
    
    problem_text = '\n'.join(problem_lines).strip()
    # Take first 1000 chars of problem description
    return problem_text[:1000] if problem_text else ""


def extract_solution_from_email(email_content: str) -> str:
    """
    Parse email content and extract the proposed solution.
    Looks for solution/recommendation/action items that are approved/final.
    Ignores alternative or rejected solutions.
    """
    if not email_content:
        return ""
    
    # Find solution-related sections
    solution_patterns = [
        r'(?:solution|proposed|recommend|action|approved|final)[:\s]+(.+?)(?=(?:alternative|option|other|instead|but|however|note|ps:|--)|$)',
    ]
    
    solution_text = ""
    for pattern in solution_patterns:
        matches = re.findall(pattern, email_content, flags=re.IGNORECASE | re.DOTALL)
        if matches:
            solution_text = matches[0].strip()
            break
    
    # If no explicit solution found, look for approved recommendations
    if not solution_text:
        # Look for "approved" or "final" solutions
        approved_match = re.search(
            r'(?:approved|final|confirmed|decided)[\s:]+(.+?)(?=(?:alternative|option|rejected|dismissed|note|ps:|--)|$)',
            email_content,
            flags=re.IGNORECASE | re.DOTALL
        )
        if approved_match:
            solution_text = approved_match.group(1).strip()
    
    # Clean up multiline solution text
    solution_text = re.sub(r'\n{2,}', ' ', solution_text)
    solution_text = re.sub(r'\s+', ' ', solution_text).strip()
    
    return solution_text[:1000] if solution_text else ""


def correlate_email_with_problem(email_content: str, user_problem: str, payload: Dict[str, Any]) -> Dict[str, str]:
    """
    Read email, correlate with user problem/solution inputs, and generate PSS output.
    The full email text is always passed to the AI so it can extract context itself.
    """
    email_text = (email_content or "").strip()
    user_text = (user_problem or "").strip()

    # Guard: if email reading itself returned an error, surface it early
    if email_text.startswith("Error reading"):
        return {
            "title": "", "problem_statement": "", "solution_statement": "",
            "raw": email_text,
            "error": email_text,
        }

    # Get user's existing proposed solution from payload
    user_solution = str(
        payload.get("current_solution_text") or payload.get("proposed_solution_text") or ""
    ).strip()

    # Always include the full email so the AI can read it directly
    input_parts = []
    if email_text:
        input_parts.append(f"Email Content:\n{email_text[:3000]}")
    if user_text:
        input_parts.append(f"User Entered Problem:\n{user_text}")
    if user_solution:
        input_parts.append(f"User Proposed Solution:\n{user_solution}")

    combined_input = "\n\n".join(input_parts)
    
    # Build payload for AI processing
    email_payload = dict(payload)
    email_payload["current_problem_text"] = combined_input
    email_payload["email_content"] = email_text
    # Ensure user's proposed solution is preserved in payload
    if user_solution and not email_payload.get("proposed_solution_text"):
        email_payload["proposed_solution_text"] = user_solution
    
    # Generate PSS from correlated content
    return generate_full_pss(email_payload)


# ----------------------------
# Optional standalone Tk action
# ----------------------------
def reframe_action():
    try:
        user_input = user_text.get("1.0", tk.END).strip()
        file_text = file_content.get("1.0", tk.END).strip()
    except Exception:
        return

    combined = ""
    if user_input:
        combined += f"User Problem:\n{user_input}\n\n"
    if file_text:
        combined += f"Reference File Content:\n{file_text}\n\n"

    if not combined.strip():
        messagebox.showwarning("Input Required", "Please provide text or browse a file.")
        return

    bom_context = build_bom_context(combined)
    if bom_context:
        combined += bom_context

    output_text.delete("1.0", tk.END)
    output_text.insert(tk.END, "Processing...\n")
    root.update_idletasks()

    reframed = reframe_problem(combined)
    output_text.delete("1.0", tk.END)
    output_text.insert(tk.END, reframed)